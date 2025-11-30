# app/tasks/document_processing_tasks.py
"""
Async document processing tasks using Celery
Handles file uploads, OCR, and text extraction
"""
import asyncio
import logging
import os
from typing import Dict, Any
from celery import Task
from celery.exceptions import SoftTimeLimitExceeded
import hashlib

from app.celery_app import celery_app
from app.services.cache_service import CacheService

logger = logging.getLogger(__name__)
cache_service = CacheService()

class DocumentProcessingTask(Task):
    """Base task with rate limiting and error handling"""
    autoretry_for = (Exception,)
    retry_kwargs = {'max_retries': 3, 'countdown': 5}
    rate_limit = '5/m'  # 5 tasks per minute per worker
    soft_time_limit = 300  # 5 minutes
    time_limit = 360  # 6 minutes

@celery_app.task(bind=True, base=DocumentProcessingTask, name='process_document_async')
def process_document_async(self,
                          task_id: str,
                          file_path: str,
                          file_type: str,
                          user_id: int = None):
    """
    Async document processing task
    Handles PDF, DOCX, images with OCR
    """
    try:
        # Update task status
        cache_service.set_task_status(task_id, {
            'status': 'processing',
            'progress': 10,
            'message': 'Starting document processing...'
        })
        
        # Check if file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Check cache first
        cache_key = f"doc:{hashlib.md5(f'{file_path}:{file_type}'.encode()).hexdigest()}"
        cached_result = cache_service.get(cache_key)
        
        if cached_result:
            logger.info(f"Returning cached document for task {task_id}")
            cache_service.set_task_status(task_id, {
                'status': 'completed',
                'progress': 100,
                'result': cached_result
            })
            return {'task_id': task_id, 'from_cache': True}
        
        # Update progress
        cache_service.set_task_status(task_id, {
            'status': 'processing',
            'progress': 30,
            'message': f'Processing {file_type} file...'
        })
        
        # Process based on file type
        extracted_text = ""
        
        if file_type in ['pdf', 'application/pdf']:
            extracted_text = process_pdf(file_path)
        elif file_type in ['docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
            extracted_text = process_docx(file_path)
        elif file_type in ['txt', 'text/plain']:
            extracted_text = process_text(file_path)
        elif file_type in ['png', 'jpg', 'jpeg', 'image/png', 'image/jpeg']:
            extracted_text = process_image_ocr(file_path)
        elif file_type in ['pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
            extracted_text = process_pptx(file_path)
        elif file_type in ['xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            extracted_text = process_xlsx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Update progress
        cache_service.set_task_status(task_id, {
            'status': 'processing',
            'progress': 80,
            'message': 'Finalizing extraction...'
        })
        
        # Prepare result
        result = {
            'text': extracted_text,
            'file_path': file_path,
            'file_type': file_type,
            'word_count': len(extracted_text.split()),
            'char_count': len(extracted_text),
            'user_id': user_id
        }
        
        # Cache result
        cache_service.set(cache_key, result, ttl=3600)
        
        # Store final result
        cache_service.set_task_status(task_id, {
            'status': 'completed',
            'progress': 100,
            'result': result
        })
        
        logger.info(f"Document processing completed for task {task_id}")
        return {'task_id': task_id, 'success': True}
        
    except SoftTimeLimitExceeded:
        logger.error(f"Task {task_id} timed out")
        cache_service.set_task_status(task_id, {
            'status': 'failed',
            'error': 'Task timed out. File might be too large.'
        })
        raise
        
    except Exception as e:
        logger.error(f"Error in document processing task {task_id}: {str(e)}")
        cache_service.set_task_status(task_id, {
            'status': 'failed',
            'error': str(e)
        })
        raise

def process_pdf(file_path: str) -> str:
    """Extract text from PDF"""
    try:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF processing error: {str(e)}")
        raise

def process_docx(file_path: str) -> str:
    """Extract text from DOCX"""
    try:
        import docx
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
        return text.strip()
    except Exception as e:
        logger.error(f"DOCX processing error: {str(e)}")
        raise

def process_text(file_path: str) -> str:
    """Read plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read().strip()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, 'r', encoding='latin-1') as file:
            return file.read().strip()
    except Exception as e:
        logger.error(f"Text processing error: {str(e)}")
        raise

def process_image_ocr(file_path: str) -> str:
    """Extract text from image using OCR"""
    try:
        import pytesseract
        from PIL import Image
        
        image = Image.open(file_path)
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        text = pytesseract.image_to_string(image, lang='eng')
        return text.strip()
    except ImportError:
        logger.warning("pytesseract not available, returning empty string")
        return ""
    except Exception as e:
        logger.error(f"OCR processing error: {str(e)}")
        return ""

def process_pptx(file_path: str) -> str:
    """Extract text from PowerPoint"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PPTX processing error: {str(e)}")
        raise

def process_xlsx(file_path: str) -> str:
    """Extract text from Excel"""
    try:
        import openpyxl
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"XLSX processing error: {str(e)}")
        raise